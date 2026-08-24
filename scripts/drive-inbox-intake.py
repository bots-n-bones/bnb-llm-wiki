#!/usr/bin/env python3
"""Incrementally stage Google Drive Inbox files for governed KB curation.

The script never mutates Drive and never publishes to the Wiki. It writes a
private staging area containing metadata, bounded text extracts, a recursive
manifest, checksums, run receipts, and tombstones.
"""

import argparse
import csv
import hashlib
import io
import json
import mimetypes
import os
import re
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path

FOLDER_MIME = "application/vnd.google-apps.folder"
GOOGLE_DOC = "application/vnd.google-apps.document"
GOOGLE_SHEET = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDES = "application/vnd.google-apps.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PDF_MIME = "application/pdf"
TEXT_MIMES = {
    "application/json", "application/xml", "application/yaml",
    "text/csv", "text/markdown", "text/plain", "text/tab-separated-values",
}
MEDIA_PREFIXES = ("audio/", "image/", "video/")
PROJECTS = ("svmpx", "hello-i-am", "content-os", "bots-n-bones", "quntm", "ursus", "hermes", "shared")
PROJECT_HINTS = {
    "svmpx": ("svmpx", "savimpex", "ddp", "source offer", "client proposal"),
    "hello-i-am": ("hello i am", "hello-iam", "helloiam", "instagram launch"),
    "content-os": ("content os", "content plan", "content calendar", "content factory", "tone and voice"),
    "bots-n-bones": ("bots-n-bones", "bots n bones", "bnb"),
    "quntm": ("quntm", "quantum"),
    "ursus": ("ursus",),
    "hermes": ("hermes", "гермес"),
}
RESTRICTED_HINTS = ("password", "secret", "credential", "private key", "api key", "token")


def utc_now():
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def safe_slug(value):
    value = re.sub(r"[^a-z0-9]+", "-", value.casefold()).strip("-")
    return value[:72] or "untitled"


def sha256_bytes(value):
    return hashlib.sha256(value).hexdigest()


def atomic_write(path, content, mode=0o600):
    path.parent.mkdir(parents=True, exist_ok=True)
    data = content.encode("utf-8") if isinstance(content, str) else content
    with tempfile.NamedTemporaryFile(prefix=f".{path.name}.", dir=path.parent, delete=False) as tmp:
        tmp.write(data)
        temporary = Path(tmp.name)
    os.chmod(temporary, mode)
    os.replace(temporary, path)


def read_json(path, default):
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (FileNotFoundError, json.JSONDecodeError):
        return default


def project_for(path, name, description=""):
    haystack = " ".join((path, name, description or "")).casefold()
    for project, hints in PROJECT_HINTS.items():
        if any(hint in haystack for hint in hints):
            return project
    return "shared"


def confidentiality_for(path, name, description=""):
    haystack = " ".join((path, name, description or "")).casefold()
    return "restricted" if any(hint in haystack for hint in RESTRICTED_HINTS) else "internal"


def fingerprint(item):
    values = [item.get(key) for key in ("id", "version", "modifiedTime", "md5Checksum", "size", "mimeType")]
    return sha256_bytes("\x1f".join("" if value is None else str(value) for value in values).encode("utf-8"))


def list_children(service, folder_id):
    token = None
    while True:
        result = service.files().list(
            q=f"'{folder_id}' in parents and trashed = false",
            pageSize=1000,
            pageToken=token,
            fields=("nextPageToken,files(id,name,mimeType,size,createdTime,modifiedTime,version,"
                    "md5Checksum,webViewLink,parents,description)"),
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()
        yield from result.get("files", [])
        token = result.get("nextPageToken")
        if not token:
            break


def walk_drive(service, root_id, root_name="00 - inbox"):
    stack = [(root_id, root_name)]
    seen = set()
    while stack:
        folder_id, folder_path = stack.pop()
        if folder_id in seen:
            continue
        seen.add(folder_id)
        children = sorted(list_children(service, folder_id), key=lambda item: (item.get("name", "").casefold(), item["id"]))
        for item in children:
            item = dict(item)
            item["parent_folder_id"] = folder_id
            item["original_path"] = f"{folder_path}/{item.get('name', item['id'])}"
            yield item
            if item.get("mimeType") == FOLDER_MIME:
                stack.append((item["id"], item["original_path"]))


def list_all_visible(service):
    """List visible Drive metadata in page-sized calls, without file downloads."""
    token = None
    while True:
        result = service.files().list(
            q="trashed = false",
            pageSize=1000,
            pageToken=token,
            fields=("nextPageToken,files(id,name,mimeType,size,createdTime,modifiedTime,version,"
                    "md5Checksum,webViewLink,parents,description)"),
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()
        yield from result.get("files", [])
        token = result.get("nextPageToken")
        if not token:
            break


def descendants_from_inventory(items, root_id, root_name="00 - inbox"):
    """Resolve one root's descendants from a complete metadata inventory."""
    children = {}
    for raw in items:
        item = dict(raw)
        for parent_id in item.get("parents") or []:
            children.setdefault(parent_id, []).append(item)

    stack = [(root_id, root_name)]
    seen = set()
    while stack:
        folder_id, folder_path = stack.pop()
        if folder_id in seen:
            continue
        seen.add(folder_id)
        entries = sorted(children.get(folder_id, []), key=lambda item: (item.get("name", "").casefold(), item["id"]))
        folders = []
        for raw in entries:
            item = dict(raw)
            item["parent_folder_id"] = folder_id
            item["original_path"] = f"{folder_path}/{item.get('name', item['id'])}"
            yield item
            if item.get("mimeType") == FOLDER_MIME:
                folders.append((item["id"], item["original_path"]))
        stack.extend(reversed(folders))


def walk_drive_bulk(service, root_id, root_name="00 - inbox"):
    return descendants_from_inventory(list_all_visible(service), root_id, root_name)


def download_bytes(service, file_id, export_mime=None, max_bytes=25_000_000):
    from googleapiclient.http import MediaIoBaseDownload
    request = service.files().export_media(fileId=file_id, mimeType=export_mime) if export_mime else service.files().get_media(fileId=file_id)
    stream = io.BytesIO()
    downloader = MediaIoBaseDownload(stream, request, chunksize=1024 * 1024)
    done = False
    while not done:
        status, done = downloader.next_chunk()
        if status and status.resumable_progress > max_bytes:
            raise ValueError(f"download exceeds {max_bytes} bytes")
    value = stream.getvalue()
    if len(value) > max_bytes:
        raise ValueError(f"download exceeds {max_bytes} bytes")
    return value


def extract_xlsx(raw):
    from openpyxl import load_workbook
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    chunks = []
    for sheet in workbook.worksheets:
        chunks.append(f"# Sheet: {sheet.title}")
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 5000:
                chunks.append("[sheet truncated after 5000 rows]")
                break
            values = ["" if value is None else str(value).replace("\n", " ") for value in row]
            if any(values):
                chunks.append("\t".join(values).rstrip())
    return "\n".join(chunks)


def extract_docx(raw):
    from docx import Document
    document = Document(io.BytesIO(raw))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text.replace("\n", " ") for cell in row.cells))
    return "\n".join(parts)


def extract_pptx(raw):
    from pptx import Presentation
    presentation = Presentation(io.BytesIO(raw))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def extract_pdf(raw):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def extract_item(service, item, max_bytes, max_chars):
    mime = item.get("mimeType", "application/octet-stream")
    if mime.startswith(MEDIA_PREFIXES):
        return "metadata-only", "", None, "media is registered without downloading"
    try:
        if mime == GOOGLE_DOC:
            raw = download_bytes(service, item["id"], "text/plain", max_bytes)
            text = raw.decode("utf-8", errors="replace")
        elif mime == GOOGLE_SHEET:
            raw = download_bytes(service, item["id"], XLSX_MIME, max_bytes)
            text = extract_xlsx(raw)
        elif mime == GOOGLE_SLIDES:
            raw = download_bytes(service, item["id"], PPTX_MIME, max_bytes)
            text = extract_pptx(raw)
        else:
            raw = download_bytes(service, item["id"], None, max_bytes)
            if mime == PDF_MIME or item.get("name", "").casefold().endswith(".pdf"):
                text = extract_pdf(raw)
            elif mime == XLSX_MIME or item.get("name", "").casefold().endswith((".xlsx", ".xlsm")):
                text = extract_xlsx(raw)
            elif mime == DOCX_MIME or item.get("name", "").casefold().endswith(".docx"):
                text = extract_docx(raw)
            elif mime == PPTX_MIME or item.get("name", "").casefold().endswith(".pptx"):
                text = extract_pptx(raw)
            elif mime in TEXT_MIMES or mime.startswith("text/"):
                text = raw.decode("utf-8", errors="replace")
            else:
                return "metadata-only", "", sha256_bytes(raw), "unsupported binary is registered without text"
        normalized = text.replace("\x00", "").strip()
        truncated = len(normalized) > max_chars
        if truncated:
            normalized = normalized[:max_chars] + "\n\n[extract truncated]"
        return "extracted", normalized, sha256_bytes(normalized.encode("utf-8")), "text truncated" if truncated else None
    except Exception as exc:
        return "extract-error", "", None, f"{type(exc).__name__}: {exc}"


def manifest_record(item, snapshot_id, captured_at, export_sha=None):
    mime = item.get("mimeType", "application/octet-stream")
    return {
        "snapshot_id": snapshot_id,
        "captured_at": captured_at,
        "drive_file_id": item["id"],
        "name": item.get("name", ""),
        "original_path": item.get("original_path"),
        "mime_type": mime,
        "url": item.get("webViewLink") or f"https://drive.google.com/open?id={item['id']}",
        "parent_ids": item.get("parents") or [item.get("parent_folder_id")],
        "size": int(item["size"]) if item.get("size") else None,
        "created_time": item.get("createdTime"),
        "modified_time": item.get("modifiedTime"),
        "drive_version": item.get("version"),
        "md5_checksum": item.get("md5Checksum"),
        "export_sha256": export_sha,
        "classification": "discovered" if mime == FOLDER_MIME else "supporting-source",
        "duplicate_group": None,
        "preferred_source_id": None,
        "notes": None,
    }


def stage_package(staging, item, captured_at, status, text, export_sha, note):
    classification_context = " ".join((item.get("description", ""), text[:20_000]))
    project = project_for(item.get("original_path", ""), item.get("name", ""), classification_context)
    package = {
        "schema": "hermes-kb-intake/v1",
        "intake_id": f"drive-{item['id']}",
        "drive_file_id": item["id"],
        "title": item.get("name", item["id"]),
        "project": project,
        "confidentiality": confidentiality_for(item.get("original_path", ""), item.get("name", ""), item.get("description", "")),
        "status": "pending-review",
        "extraction_status": status,
        "extraction_note": note,
        "extracted_text": text,
        "extract_sha256": export_sha,
        "captured_at": captured_at,
        "fingerprint": fingerprint(item),
        "source": manifest_record(item, "pending", captured_at, export_sha),
        "publication_approved_at": None,
        "publication_approved_by": None,
        "materialized_at": None,
    }
    path = staging / "pending" / f"{safe_slug(item.get('name', 'source'))}-{item['id'][:12]}.json"
    atomic_write(path, json.dumps(package, ensure_ascii=False, indent=2) + "\n")
    return path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--folder-id", required=True)
    parser.add_argument("--folder-name", default="00 - inbox")
    parser.add_argument("--token", required=True)
    parser.add_argument("--staging", required=True)
    parser.add_argument("--max-bytes", type=int, default=25_000_000)
    parser.add_argument("--max-chars", type=int, default=500_000)
    parser.add_argument("--manifest-only", action="store_true")
    parser.add_argument("--bulk-list", action="store_true", help="build the root tree from one paginated metadata inventory")
    args = parser.parse_args()

    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build

    staging = Path(args.staging).resolve()
    staging.mkdir(parents=True, exist_ok=True)
    captured_at = utc_now()
    snapshot_id = "drive-inbox-" + captured_at.replace(":", "").replace("-", "")
    state_path = staging / "state.json"
    state = read_json(state_path, {"files": {}})
    previous = state.get("files", {})
    current = {}
    records = []
    changed = []
    errors = []

    credentials = Credentials.from_authorized_user_file(args.token, ["https://www.googleapis.com/auth/drive.readonly"])
    service = build("drive", "v3", credentials=credentials, cache_discovery=False)
    walker = walk_drive_bulk if args.bulk_list else walk_drive
    items = list(walker(service, args.folder_id, args.folder_name))
    for item in items:
        fp = fingerprint(item)
        current[item["id"]] = {"fingerprint": fp, "path": item.get("original_path"), "mime_type": item.get("mimeType")}
        if item.get("mimeType") == FOLDER_MIME:
            records.append(manifest_record(item, snapshot_id, captured_at))
            continue
        if args.manifest_only:
            current[item["id"]]["extraction_status"] = "metadata-only"
            records.append(manifest_record(item, snapshot_id, captured_at))
            continue
        prior = previous.get(item["id"], {})
        if prior.get("fingerprint") == fp and prior.get("extraction_status") in {"extracted", "metadata-only"}:
            current[item["id"]]["extraction_status"] = prior["extraction_status"]
            records.append(manifest_record(item, snapshot_id, captured_at, prior.get("export_sha256")))
            continue
        status, text, export_sha, note = extract_item(service, item, args.max_bytes, args.max_chars)
        package_path = stage_package(staging, item, captured_at, status, text, export_sha, note)
        current[item["id"]].update({"extraction_status": status, "export_sha256": export_sha, "staged_path": str(package_path)})
        records.append(manifest_record(item, snapshot_id, captured_at, export_sha))
        changed.append({"id": item["id"], "path": item.get("original_path"), "status": status, "package": str(package_path)})
        if status == "extract-error":
            errors.append(changed[-1] | {"error": note})

    tombstones = [{"id": file_id, **metadata} for file_id, metadata in previous.items() if file_id not in current]
    manifest_text = "".join(json.dumps(record, ensure_ascii=False, sort_keys=True) + "\n" for record in records)
    stable_records = []
    for record in records:
        stable = dict(record)
        stable.pop("snapshot_id", None)
        stable.pop("captured_at", None)
        stable_records.append(stable)
    inventory_text = "".join(json.dumps(record, ensure_ascii=False, sort_keys=True) + "\n" for record in stable_records)
    inventory_sha256 = sha256_bytes(inventory_text.encode("utf-8"))
    manifest_path = staging / "manifest.jsonl"
    atomic_write(manifest_path, manifest_text)
    atomic_write(staging / "manifest.sha256", sha256_bytes(manifest_text.encode("utf-8")) + "  manifest.jsonl\n")
    atomic_write(staging / "inventory.sha256", inventory_sha256 + "  stable-inventory\n")
    csv_stream = io.StringIO()
    csv_fields = ["snapshot_id", "captured_at", "drive_file_id", "name", "original_path", "mime_type", "url", "parent_ids", "size", "created_time", "modified_time", "drive_version", "md5_checksum", "export_sha256", "classification"]
    writer = csv.DictWriter(csv_stream, fieldnames=csv_fields, extrasaction="ignore")
    writer.writeheader()
    for record in records:
        row = dict(record)
        row["parent_ids"] = json.dumps(row.get("parent_ids") or [], ensure_ascii=False)
        writer.writerow(row)
    csv_text = csv_stream.getvalue()
    atomic_write(staging / "manifest.csv", csv_text)
    atomic_write(staging / "manifest.csv.sha256", sha256_bytes(csv_text.encode("utf-8")) + "  manifest.csv\n")
    new_state = {"schema": "hermes-kb-drive-state/v1", "folder_id": args.folder_id, "captured_at": captured_at, "files": current}
    atomic_write(state_path, json.dumps(new_state, ensure_ascii=False, indent=2) + "\n")
    receipt = {
        "schema": "hermes-kb-intake-run/v1", "snapshot_id": snapshot_id, "captured_at": captured_at,
        "folder_id": args.folder_id, "item_count": len(items), "file_count": sum(1 for item in items if item.get("mimeType") != FOLDER_MIME),
        "changed_count": len(changed), "error_count": len(errors), "tombstone_count": len(tombstones),
        "changed": changed, "errors": errors, "tombstones": tombstones,
        "manifest_sha256": sha256_bytes(manifest_text.encode("utf-8")), "inventory_sha256": inventory_sha256,
    }
    atomic_write(staging / "runs" / f"{snapshot_id}.json", json.dumps(receipt, ensure_ascii=False, indent=2) + "\n")
    atomic_write(staging / "health.json", json.dumps(receipt, ensure_ascii=False, indent=2) + "\n", mode=0o644)
    print(json.dumps(receipt, ensure_ascii=False))
    return 1 if errors else 0


if __name__ == "__main__":
    sys.exit(main())
